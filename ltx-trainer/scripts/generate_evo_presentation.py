"""
Generate a PowerPoint presentation comparing Evolution Strategies vs Gradient Descent
for neural network fine-tuning, specifically in the context of SCD Evolution for LTX-2.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Color palette ──────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
BG_CARD      = RGBColor(0x1A, 0x2B, 0x3D)   # Slightly lighter navy
ACCENT_BLUE  = RGBColor(0x00, 0xB4, 0xD8)   # Cyan-blue
ACCENT_GREEN = RGBColor(0x2D, 0xC6, 0x5D)   # Green
ACCENT_RED   = RGBColor(0xFF, 0x5C, 0x5C)   # Red/orange
ACCENT_GOLD  = RGBColor(0xFF, 0xD1, 0x66)   # Gold/yellow
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE0)
MID_GRAY     = RGBColor(0x6B, 0x8A, 0xA8)
ROW_ALT      = RGBColor(0x14, 0x22, 0x33)   # Alternating table row

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(18), bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a colored title bar at the top of the slide."""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BG_CARD)
    # Accent line
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Pt(3), fill_color=ACCENT_BLUE)
    # Title text
    add_textbox(slide, title_text,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7),
                font_size=Pt(28), bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.4), Inches(0.8), Inches(12.5), Inches(0.35),
                    font_size=Pt(14), color=ACCENT_BLUE)
    return bar


def add_slide_number(slide, num, total):
    add_textbox(slide, f"{num} / {total}",
                Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
                font_size=Pt(11), color=MID_GRAY, align=PP_ALIGN.RIGHT)


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=Pt(17), bullet_color=ACCENT_BLUE,
                    text_color=LIGHT_GRAY, indent_levels=None):
    """Add a bulleted list. items is list of strings or (text, level) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = level
        indent = "    " * level
        bullet = "•" if level == 0 else "◦"
        run = p.add_run()
        run.text = f"{indent}{bullet}  {text}"
        run.font.size = font_size if level == 0 else Pt(font_size.pt - 2)
        run.font.color.rgb = text_color if level > 0 else WHITE
        p.space_before = Pt(6 if level == 0 else 3)
        p.space_after = Pt(2)
    return txBox


def add_table(slide, headers, rows, left, top, width, col_widths=None,
              header_fill=ACCENT_BLUE, header_text=WHITE,
              row_fill=BG_CARD, row_alt_fill=ROW_ALT,
              row_text=LIGHT_GRAY, font_size=Pt(15)):
    """Add a formatted table to the slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    height = Inches(0.45) * num_rows

    tbl = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    else:
        col_w = width // num_cols
        for i in range(num_cols):
            tbl.columns[i].width = col_w

    def set_cell(cell, text, fill_color, text_color, bold=False, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = text_color
        # Remove cell border lines
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

    # Header row
    for j, h in enumerate(headers):
        set_cell(tbl.cell(0, j), h, header_fill, header_text, bold=True, align=PP_ALIGN.CENTER)

    # Data rows
    for i, row in enumerate(rows):
        fill = row_fill if i % 2 == 0 else row_alt_fill
        for j, val in enumerate(row):
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            set_cell(tbl.cell(i + 1, j), val, fill, row_text, align=align)

    return tbl


def add_highlight_box(slide, text, left, top, width, height,
                      bg_color=BG_CARD, border_color=ACCENT_BLUE,
                      text_color=WHITE, font_size=Pt(16), bold=False, italic=False):
    box = add_rect(slide, left, top, width, height,
                   fill_color=bg_color, line_color=border_color, line_width=Pt(1.5))
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = text_color
    return box


# ══════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # Completely blank layout
TOTAL_SLIDES = 10


# ── Slide 1: Title ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

# Decorative accent bar left
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=ACCENT_BLUE)

# Big title
add_textbox(slide, "Evolution Strategies vs Gradient Descent",
            Inches(0.6), Inches(1.6), Inches(12.0), Inches(1.2),
            font_size=Pt(42), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle line 1
add_textbox(slide, "For Autoregressive Video Quality Optimization",
            Inches(0.6), Inches(2.9), Inches(12.0), Inches(0.7),
            font_size=Pt(26), bold=False, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(3.0), Inches(3.7), Inches(7.3), Pt(2), fill_color=ACCENT_GOLD)

# Subtitle line 2
add_textbox(slide, "SCD Evolution  —  LTX-2 19B",
            Inches(0.6), Inches(3.85), Inches(12.0), Inches(0.6),
            font_size=Pt(22), bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Bottom tagline
add_textbox(slide, "Optimizing Non-Differentiable Autoregressive Objectives via Gradient-Free Search",
            Inches(0.6), Inches(5.0), Inches(12.0), Inches(0.5),
            font_size=Pt(15), italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)


# ── Slide 2: The Core Problem ──────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "The Core Problem", "Why gradient descent fails for autoregressive video")
add_slide_number(slide, 2, TOTAL_SLIDES)

bullets = [
    "Teacher forcing trains on CLEAN frames, but inference runs AUTOREGRESSIVELY",
    ("Model never sees its own errors during training → optimistic loss landscape", 1),
    "Compounding errors invisible during training → quality collapse at inference",
    ("Frame 1 error feeds into Frame 2 → Frame 3 → exponential degradation", 1),
    "GD optimizes a PROXY loss (per-step MSE); ES evaluates the ACTUAL inference pipeline",
    ("If the objective is non-differentiable, GD simply cannot be used", 1),
]
add_bullet_list(slide, bullets, Inches(0.4), Inches(1.35), Inches(8.5), Inches(4.2),
                font_size=Pt(17))

# Equation box
add_highlight_box(
    slide,
    "fitness  =  AR_rollout(θ)   ←  non-differentiable black box\n\n"
    "        ∂ fitness / ∂θ   =   undefined",
    Inches(9.0), Inches(1.4), Inches(4.0), Inches(2.0),
    bg_color=RGBColor(0x08, 0x12, 0x20),
    border_color=ACCENT_GOLD,
    text_color=ACCENT_GOLD,
    font_size=Pt(14),
    bold=True
)

# Key insight box
add_highlight_box(
    slide,
    "Key Insight:\nES samples from θ + σε directly in weight space, evaluates the full "
    "autoregressive rollout, and estimates the gradient via finite differences — "
    "never computing ∂L/∂θ at all.",
    Inches(9.0), Inches(3.6), Inches(4.0), Inches(2.8),
    bg_color=BG_CARD,
    border_color=ACCENT_GREEN,
    text_color=LIGHT_GRAY,
    font_size=Pt(13)
)


# ── Slide 3: FLOPs Comparison ──────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "FLOPs Comparison", "Per-update compute budget breakdown")
add_slide_number(slide, 3, TOTAL_SLIDES)

headers = ["Property", "Gradient Descent", "Evolution Strategies"]
rows = [
    ("Forward passes / update",  "1",                      "2 × pop_size  (32)"),
    ("Backward pass cost",        "~2× forward",            "None"),
    ("Total FLOPs / update",      "~3 fwd equiv",           "32 fwd"),
    ("Works on quantized models", "No  (autograd breaks)",  "Yes  ✓"),
    ("Works on non-diff rewards", "No",                     "Yes  ✓"),
    ("Memory per update",         "Grads + optimizer state","Inference only"),
    ("Parameter coverage",        "All trainable layers",   "Any subset (LoRA)"),
    ("Parallelism",               "Data parallel only",     "Embarrassingly parallel"),
]

col_widths = [Inches(3.6), Inches(4.3), Inches(4.3)]
add_table(slide, headers, rows,
          Inches(0.55), Inches(1.35), Inches(12.2),
          col_widths=col_widths, font_size=Pt(14))

note = ("Note: ES uses 32× more forward passes per update but stores zero gradients. "
        "At int8-quanto, 32 forward passes fit in ~22GB — GD backward would require ~40GB+.")
add_textbox(slide, note,
            Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.9),
            font_size=Pt(13), italic=True, color=MID_GRAY)


# ── Slide 4: When ES Wins ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "When ES Wins", "Evolution Strategies outperforms gradient descent")
add_slide_number(slide, 4, TOTAL_SLIDES)

wins = [
    "Quantized models  (int8 / fp8 breaks autograd — ufunc_add not implemented for Float8_e4m3fn)",
    "Non-differentiable fitness  (LPIPS, SSIM, VMAF, perceptual quality, AR rollout score)",
    "Memory-constrained hardware  (no backward graph stored — 12× reduction vs BPTT)",
    "Embarrassingly parallel workloads  (each candidate is independent — scalar reward communication only)",
    "Hyperparameter-robust search  (no learning rate, momentum, or scheduler tuning required)",
]
add_bullet_list(slide, wins, Inches(0.4), Inches(1.35), Inches(8.6), Inches(3.8),
                font_size=Pt(17))

# Historical callout
add_highlight_box(
    slide,
    "OpenAI (2017):\nHumanoid walking in 10 min\non 1,440 CPUs — 2,000 workers\nsharing only a scalar reward signal.\n\n"
    "MeZO (NeurIPS 2023):\nFine-tunes 30B LLM with just\ntwo forward passes — 12× less\nmemory than Adam.",
    Inches(9.1), Inches(1.4), Inches(3.9), Inches(4.2),
    bg_color=BG_CARD,
    border_color=ACCENT_GOLD,
    text_color=LIGHT_GRAY,
    font_size=Pt(14)
)

add_textbox(slide, "Our case: All four ES advantages apply simultaneously.",
            Inches(0.4), Inches(5.3), Inches(8.5), Inches(0.5),
            font_size=Pt(15), bold=True, color=ACCENT_GREEN)


# ── Slide 5: When GD Wins ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "When GD Wins", "Gradient descent remains superior for these use cases")
add_slide_number(slide, 5, TOTAL_SLIDES)

gd_wins = [
    "Smooth differentiable objectives  (MSE, cross-entropy, reconstruction losses)",
    ("GD is ~1,000× more sample-efficient when gradients are well-defined — OpenAI 2017", 1),
    "Small parameter spaces where gradient information is dense and meaningful",
    "Standard fine-tuning tasks  (text, image classification, regression)",
    ("Adam + LoRA is sufficient for 95% of fine-tuning scenarios", 1),
    "When teacher-forcing loss IS the actual deployment objective (no AR gap)",
]
add_bullet_list(slide, gd_wins, Inches(0.4), Inches(1.35), Inches(8.8), Inches(4.0),
                font_size=Pt(17))

# Verdict box
add_highlight_box(
    slide,
    "Verdict:\nGD wins on standard tasks.\nES wins when the evaluation\nfunction is a black box,\nnon-differentiable, or requires\nrunning the actual deployment\npipeline end-to-end.",
    Inches(9.1), Inches(1.4), Inches(3.9), Inches(3.8),
    bg_color=BG_CARD,
    border_color=ACCENT_RED,
    text_color=LIGHT_GRAY,
    font_size=Pt(15)
)

add_textbox(slide, "Both are tools — the right choice depends on the objective, not dogma.",
            Inches(0.4), Inches(5.5), Inches(12.0), Inches(0.5),
            font_size=Pt(15), bold=True, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)


# ── Slide 6: SCD Evolution Architecture ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "SCD Evolution Architecture", "How the fitness pipeline operates on LTX-2 19B")
add_slide_number(slide, 6, TOTAL_SLIDES)

# Pipeline description
add_highlight_box(
    slide,
    "Pipeline:  θ  →  perturb (θ + σε, θ − σε)  →  θ'  →  [8 steps × 3 AR frames × CFG]  →  fitness",
    Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.65),
    bg_color=RGBColor(0x08, 0x12, 0x20),
    border_color=ACCENT_BLUE,
    text_color=ACCENT_BLUE,
    font_size=Pt(16),
    bold=True
)

compute_items = [
    "Per generation compute:",
    ("16 antithetic pairs  ×  2 (pos/neg)  ×  6 prompts  ×  3 frames  ×  8 steps  ×  2 (CFG)  =  9,216 decoder forwards", 1),
    ("GPU-batched to ~768 calls at batch_size=12 each — all on RTX 5090 (cuda:0)", 1),
    "Memory profile:",
    ("~22GB VRAM  (int8-quanto transformer)  +  VAE on RTX PRO 4000  (cuda:1)", 1),
    ("Zero gradient storage — inference_mode() throughout", 1),
    "Fitness function: LPIPS perceptual distance + temporal consistency + optional SSIM",
    ("Lower LPIPS = better perceptual quality vs ground-truth target frames", 1),
    "Update rule: θ ← θ + α · (1/2nσ) · Σ F(θ + σεᵢ) · εᵢ  (antithetic ES gradient estimate)",
]
add_bullet_list(slide, compute_items, Inches(0.4), Inches(2.1), Inches(12.5), Inches(4.6),
                font_size=Pt(15))


# ── Slide 7: Training Speed Comparison ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "Training Speed Comparison", "Wall-clock time vs optimization target")
add_slide_number(slide, 7, TOTAL_SLIDES)

headers = ["Method", "Time / Update", "Total Budget", "Optimizes AR?", "Quantization OK?"]
rows = [
    ("ES  (16 pairs, int8)",      "~250s / gen",      "~21h  (300 gens)",    "Yes — directly",    "Yes  ✓"),
    ("GD teacher-forcing",        "~8s / step",       "~4.4h  (2000 steps)", "No — proxy loss",   "No"),
    ("GD BPTT  (8 AR steps)",     "Not feasible",     "OOM on 32GB",         "Partial",           "No"),
    ("Self Forcing (NeurIPS'25)", "~150s / step",     "~83h",                "Yes",               "No"),
    ("MeZO (LLM analogy)",        "2 fwd passes",     "Competitive",         "With AR reward",    "Yes  ✓"),
]

col_widths = [Inches(2.8), Inches(2.1), Inches(2.3), Inches(2.4), Inches(2.4)]
add_table(slide, headers, rows,
          Inches(0.3), Inches(1.35), Inches(12.6),
          col_widths=col_widths, font_size=Pt(13))

takeaway = ("ES is 2.7× slower per-update but is the ONLY method that can optimize "
            "non-differentiable AR rollout quality on quantized models within 32GB VRAM.")
add_highlight_box(slide, takeaway,
                  Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.8),
                  bg_color=BG_CARD, border_color=ACCENT_GREEN,
                  text_color=LIGHT_GRAY, font_size=Pt(15))


# ── Slide 8: The RLHF Analogy ─────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "The RLHF Analogy", "ES for video generation mirrors RLHF for language models")
add_slide_number(slide, 8, TOTAL_SLIDES)

# Two columns side by side
# LLM side
add_highlight_box(slide, "RLHF  (LLMs)",
                  Inches(0.4), Inches(1.4), Inches(5.9), Inches(0.5),
                  bg_color=RGBColor(0x00, 0x50, 0x80),
                  border_color=ACCENT_BLUE,
                  text_color=WHITE, font_size=Pt(16), bold=True)

llm_items = [
    "LLM generates tokens autoregressively",
    ("Exposure bias: training ≠ inference distribution", 1),
    "MLE/SFT creates teacher-forcing gap",
    "RLHF aligns via non-differentiable human reward",
    ("PPO / GRPO / DPO estimate policy gradient", 1),
    "Result: ChatGPT, Claude, Gemini",
]
add_bullet_list(slide, llm_items,
                Inches(0.4), Inches(2.0), Inches(5.9), Inches(3.8),
                font_size=Pt(15))

# Video side
add_highlight_box(slide, "ES Evolution  (Video Diffusion)",
                  Inches(6.9), Inches(1.4), Inches(6.0), Inches(0.5),
                  bg_color=RGBColor(0x00, 0x55, 0x30),
                  border_color=ACCENT_GREEN,
                  text_color=WHITE, font_size=Pt(16), bold=True)

vid_items = [
    "SCD generates frames autoregressively",
    ("Exposure bias: teacher-forcing ≠ AR inference", 1),
    "MSE loss creates training-inference gap",
    "ES aligns via non-differentiable perceptual fitness",
    ("Antithetic ES estimates weight-space gradient", 1),
    "Target: Long-form LTX-2 video quality",
]
add_bullet_list(slide, vid_items,
                Inches(6.9), Inches(2.0), Inches(6.0), Inches(3.8),
                font_size=Pt(15))

# Arrow between them
add_rect(slide, Inches(6.3), Inches(2.9), Inches(0.55), Inches(0.4),
         fill_color=ACCENT_GOLD)
add_textbox(slide, "≡", Inches(6.3), Inches(2.75), Inches(0.6), Inches(0.5),
            font_size=Pt(24), bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

add_textbox(slide,
            "ES for video generation  =  RLHF for text generation",
            Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.5),
            font_size=Pt(17), bold=True, italic=True,
            color=ACCENT_GOLD, align=PP_ALIGN.CENTER)


# ── Slide 9: Key Numbers ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "Key Numbers from Our Runs", "SCD Evolution on LTX-2 19B — RTX 5090 32GB")
add_slide_number(slide, 9, TOTAL_SLIDES)

# Left column: config
add_highlight_box(slide, "Configuration",
                  Inches(0.4), Inches(1.35), Inches(5.8), Inches(0.45),
                  bg_color=ACCENT_BLUE, border_color=ACCENT_BLUE,
                  text_color=WHITE, font_size=Pt(15), bold=True)

config_items = [
    "71M evolved LoRA parameters  (decoder blocks 32–47)",
    "731 training samples  (evolution_merged dataset)",
    "16 antithetic pairs per generation  (32 candidates)",
    "noise_scale: 0.02   update_scale: 0.003",
    "Fitness: LPIPS perceptual distance  (lower = better)",
    "Baseline fitness: −0.70  (normalized, pre-evolution)",
    "Quantization: int8-quanto  (~22GB VRAM)",
    "VAE: RTX PRO 4000  (cuda:1, 24GB)",
]
add_bullet_list(slide, config_items,
                Inches(0.4), Inches(1.9), Inches(5.8), Inches(4.6),
                font_size=Pt(15))

# Right column: performance
add_highlight_box(slide, "Per-Generation Compute",
                  Inches(6.6), Inches(1.35), Inches(6.3), Inches(0.45),
                  bg_color=ACCENT_GOLD, border_color=ACCENT_GOLD,
                  text_color=BG_DARK, font_size=Pt(15), bold=True)

perf_items = [
    "~250s per generation on RTX 5090",
    "~21 hours total for 300 generations",
    "9,216 decoder forwards per generation",
    "768 batched GPU calls  (batch_size=12)",
    "0 GB gradient storage  (inference_mode)",
    "~3 AR frames evaluated per candidate",
    "8 denoising steps per AR frame",
    "CFG scale: 3.0  (pos + neg = 2× cost)",
]
add_bullet_list(slide, perf_items,
                Inches(6.6), Inches(1.9), Inches(6.3), Inches(4.6),
                font_size=Pt(15))


# ── Slide 10: Key References ───────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_title_bar(slide, "Key References", "Foundational work informing SCD Evolution")
add_slide_number(slide, 10, TOTAL_SLIDES)

refs = [
    ("Evolution Strategies as a Scalable Alternative to Reinforcement Learning",
     "OpenAI, 2017  —  Salimans et al.  |  Humanoid walking in 10 min on 1,440 CPUs; "
     "1,000× less sample-efficient than GD on smooth objectives but massively parallel"),
    ("MeZO: Fine-Tuning Language Models with Just Forward Passes",
     "NeurIPS 2023  —  Zhang et al.  |  12× memory reduction over Adam; fine-tunes OPT-30B "
     "in-place via in-place SPSA; first practical LLM fine-tuning with zero backward pass"),
    ("Self Forcing: Bridging the Train-Test Gap in Autoregressive Video Diffusion",
     "NeurIPS 2025  —  VBench score 84.31  |  Trains on model's own AR predictions; "
     "requires full bf16 model + 150s/step — infeasible on quantized 32GB setup"),
    ("ES at Scale: Fine-Tuning Large Language Models via Evolution Strategies",
     "2025  —  14.4% improvement on downstream tasks with 20% of sample budget of GD"),
    ("QZO: Quantized Zeroth-Order Fine-Tuning of LLMs",
     "2025  —  18× memory reduction; ES fine-tuning of 4-bit quantized LLMs — "
     "directly analogous to our int8-quanto ES approach for video diffusion"),
]

y_pos = Inches(1.4)
for i, (title, desc) in enumerate(refs):
    color = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_GOLD, ACCENT_RED, ACCENT_BLUE][i]
    add_highlight_box(slide, f"[{i+1}] {title}\n{desc}",
                      Inches(0.4), y_pos, Inches(12.5), Inches(0.95),
                      bg_color=BG_CARD if i % 2 == 0 else ROW_ALT,
                      border_color=color,
                      text_color=LIGHT_GRAY,
                      font_size=Pt(12))
    y_pos += Inches(1.0)


# ── Save ───────────────────────────────────────────────────────────────────────
output_path = "/media/2TB/omnitransfer/inference/evo_vs_gradient_descent.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")

import os
size = os.path.getsize(output_path)
print(f"File size: {size:,} bytes  ({size/1024:.1f} KB)")
